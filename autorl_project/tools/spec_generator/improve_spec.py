
from pathlib import Path
import argparse
import json
import re
from datetime import datetime
from typing import Dict, List, Tuple, Optional

# External libs
try:
    import markdown2
    from jinja2 import Template
    from docx import Document
    from pptx import Presentation
    from pptx.util import Inches, Pt
except Exception as e:
    raise RuntimeError(
        "Missing dependency. pip install python-docx markdown2 jinja2 python-pptx"
    ) from e


# -------------------------
# Helpers: parsing & IO
# -------------------------
def read_input(path: Path) -> str:
    if path.suffix.lower() in (".md", ".markdown", ".txt"):
        return path.read_text(encoding="utf-8")
    elif path.suffix.lower() in (".docx",):
        doc = Document(str(path))
        paragraphs = []
        for p in doc.paragraphs:
            paragraphs.append(p.text)
        return "\n\n".join(paragraphs)
    else:
        raise ValueError("Unsupported input format. Use .md or .docx")


def split_top_level_sections(md: str) -> List[Tuple[str, str]]:
    """
    Very simple splitter: looks for lines starting with '## ' or '# ' as delimiters.
    Returns list of (heading, body) in order found. Heading may be empty for preamble.
    """
    lines = md.splitlines()
    sections = []
    cur_heading = ""
    cur_body_lines = []
    heading_re = re.compile(r"^\s{0,3}#{1,3}\s+(.*)")
    for line in lines:
        m = heading_re.match(line)
        if m:
            # New heading
            if cur_heading or cur_body_lines:
                sections.append((cur_heading.strip(), "\n".join(cur_body_lines).strip()))
            cur_heading = m.group(1).strip()
            cur_body_lines = []
        else:
            cur_body_lines.append(line)
    # append final
    if cur_heading or cur_body_lines:
        sections.append((cur_heading.strip(), "\n".join(cur_body_lines).strip()))
    return sections


# -------------------------
# Content templates (Jinja)
# -------------------------
ARCHITECTURE_MERMAID = """```mermaid
flowchart LR
    subgraph Mobile Device
        A[Mobile Client<br/>(screen capture, UI hooks)]
        B[On-Device Perception<br/>(OCR, Detectors, Embeddings)]
        C[Local Memory<br/>(Plan Cache, Embeddings)]
    end
    subgraph Backend
        D[Orchestrator API<br/>(FastAPI)]
        E[Planning LLM<br/>(LLM + Prompting Layer)]
        F[RL Trainer & Replay DB<br/>(Gym-like envs)]
        G[Vector DB & Plan Repo<br/>(Qdrant / SQLite)]
        H[HITL Service & Audit Logs]
        I[Monitoring & Metrics<br/>(Prometheus / Grafana)]
    end
    A --> B --> D
    D --> E
    E --> D
    D --> F
    F --> G
    D --> H
    H --> D
    G --> E
    I --> D
```"""

SECTION_ENHANCEMENT_TEMPLATE = """
### {{ title }}

{{ original | default('') }}

#### Enhanced: Key Goals
{{ goals }}

#### Architecture Details
{{ arch }}

#### Implementation Notes
{{ impl }}

#### Risk & Mitigation
| Risk | Mitigation |
|---|---|
{% for r,m in risks %}
| {{ r }} | {{ m }} |
{% endfor %}

#### Acceptance Criteria (measurable)
{% for crit in acceptance %}
- {{ crit }}
{% endfor %}
"""

CHECKLIST_TEMPLATE = """
# Hackathon Submission Checklist (Auto-generated)

- [ ] One-sentence project summary (top of README)
- [ ] Clear demo plan + 2–3 short recorded clips (phone + screen capture)
- [ ] GitHub repo with README, how to run, and sample data
- [ ] Tech spec document (this file)
- [ ] Architecture diagram (mermaid)
- [ ] Minimal runnable demo (emulator or recorded video)
- [ ] Security & privacy section (PII masking)
- [ ] Human-in-the-loop demo + explanation
- [ ] Test plan (unit + integration)
- [ ] Metrics to show (task success rate, latency, failure modes)
- [ ] Deployment instructions / Docker compose
- [ ] Slides (3–7 slides) for judges
- [ ] Short executive summary (one paragraph)
"""

SLIDE_TITLES = [
    "AutoRL — Executive Summary",
    "Problem & Opportunity",
    "How It Works — Architecture",
    "Demo Plan",
    "Technical Novelty & Innovations",
    "Evaluation Metrics & Results",
    "Roadmap & Ask"
]


# -------------------------
# Enrichment logic
# -------------------------
def enhance_section(title: str, body: str) -> str:
    """
    Return an enriched markdown section for a given title/body.
    This is deterministic and rule-based — uses heuristics tuned for AutoRL.
    """
    # Basic heuristics for content generation
    title_lower = title.lower()
    goals = ""
    arch = ""
    impl = ""
    risks = []
    acceptance = []

    if "executive" in title_lower or "summary" in title_lower:
        goals = "Convey the value prop: visual, cross-app mobile automation using vision + LLM + RL. Target judges and a non-technical demo audience."
        arch = "Mobile client captures screenshots, sends context to Orchestrator; Orchestrator calls Planning LLM to emit action plans; Execution agent runs actions via Appium/Accessibility APIs; RL training loop uses emulator replay to improve policy."
        impl = "Provide a one-paragraph demo script and a minimal command line to run demo. Include a sample dataset and expected outputs."
        risks = [("Vague value proposition", "Make the one-liner precise and include target users + metric to improve."),
                 ("No demo", "Ship a simple end-to-end recorded demo or runnable emulator script.")]
        acceptance = ["Single-sentence value proposition present", "One demo link or video included"]
    elif "architecture" in title_lower or "system architecture" in title_lower or "how it works" in title_lower:
        goals = "Show component interactions and dataflow; include mermaid diagram; explain failure modes and retries."
        arch = "See architecture diagram below.\n\n" + ARCHITECTURE_MERMAID
        impl = "Add explicit wire format examples (JSON for UI state, sample plan JSON with TAP/TYPE/SWIPE). Provide example API endpoints (POST /instruction, GET /status)."
        risks = [("Incomplete dataflow", "Add sequence diagram for critical flow (e.g., payment approval)."),
                 ("No example payloads", "Add at least 2 concrete JSON payload examples.")]
        acceptance = ["Diagram present", "Example API payloads included", "Error/retry flows documented"]
    elif "rl" in title_lower or "learning" in title_lower or "training" in title_lower:
        goals = "Describe RL environment (observations, actions, rewards), bootstrapping via imitation learning, and offline training pipeline."
        arch = "Use Gym-like AutoRLEnv with observations: (screenshot embedding, UI element list). Actions: (element_id, action_type, text). Reward shaping: +1 on task success, -0.1 per invalid action, large negative for destructive mistakes."
        impl = "List selected RL algorithms (PPO + Behavior Cloning for warmstart), training budget, checkpoints, and model evaluation protocol."
        risks = [("Sparse reward", "Use dense shaping or curriculum learning, bootstrap with imitation traces."),
                 ("Overfitting to specific app layouts", "Data augmentation and multi-app training.")]
        acceptance = ["Formal RL env spec", "Training hyperparameters", "Evaluation suite defined"]
    elif "security" in title_lower or "privacy" in title_lower:
        goals = "Document PII handling, encryption at rest/in transit, and audit logs for HITL."
        arch = "Mask or blur PII in screenshots before logging. Use role-based access for override endpoints. Audit log with append-only storage."
        impl = "Provide code snippets for masking screenshots and an example of storing audit logs."
        risks = [("PII leakage", "Ensure screenshots are masked on-device before transit."),
                 ("Weak ACLs", "Use RBAC and logging for override endpoints.")]
        acceptance = ["PII handling explained", "Audit trail present", "Override gate described"]
    elif "testing" in title_lower or "evaluation" in title_lower:
        goals = "Provide test plans and evaluation metrics tied to the judging criteria."
        arch = "Integration tests for end-to-end: instruction->plan->execution on emulator. Unit tests for perception and planner modules."
        impl = "Provide pytest examples and synthetic testcases. Add a test harness that replays recorded sessions."
        risks = [("No tests", "Add smoke tests and a CI job to run them on PRs."),
                 ("No eval metrics", "Include task success rate, latency, and regression checks.")]
        acceptance = ["Test suite included", "CI runs smoke tests", "Benchmarks included"]
    else:
        # generic section enhancement
        goals = "Clarify purpose and measurable deliverables for this section."
        arch = "If this is an actionable section, add examples & acceptance criteria."
        impl = "Add code samples, small diagrams, and TODO markers for missing details."
        risks = [("Underspecified", "Add examples and metrics")]
        acceptance = ["Clear purpose statement", "At least one measurable acceptance criterion"]

    template = Template(SECTION_ENHANCEMENT_TEMPLATE)
    return template.render(
        title=title or "Preamble",
        original=body,
        goals=goals,
        arch=arch,
        impl=impl,
        risks=risks,
        acceptance=acceptance,
    )


# -------------------------
# PPTX generator
# -------------------------
def create_slides(out_path: Path, title: str, author: str = "AutoRL Team"):
    prs = Presentation()
    # Basic title slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    subtitle_placeholder = slide.placeholders[1]
    title_placeholder.text = title
    subtitle_placeholder.text = f"Generated {datetime.utcnow().isoformat()} • {author}"

    # Add specified slides
    for t in SLIDE_TITLES:
        slide_layout = prs.slide_layouts[1]  # title + content
        s = prs.slides.add_slide(slide_layout)
        s.shapes.title.text = t
        body = s.shapes.placeholders[1].text_frame
        body.text = "- Key bullet 1"
        p = body.add_paragraph()
        p.text = "- Key bullet 2"
        p.level = 1

    prs.save(str(out_path))


# -------------------------
# Main pipeline
# -------------------------
def build_enhanced_spec(input_text: str) -> Tuple[str, Dict]:
    """
    Returns enhanced markdown string and metadata dict.
    """
    sections = split_top_level_sections(input_text)
    enhanced_sections = []
    metadata = {"generated_at": datetime.utcnow().isoformat(), "sections": []}

    for title, body in sections:
        enriched_md = enhance_section(title, body)
        enhanced_sections.append(enriched_md)
        metadata["sections"].append({"title": title or "Preamble", "has_body": bool(body)})

    # assemble final doc
    header = f"# {metadata.get('project_name','AutoRL')} — Technical Specification (Enhanced)\n\n"
    header += f"_Generated: {metadata['generated_at']}_\n\n"
    # optionally include a short executive summary region by combining first section if present
    final_md = header + "\n\n".join(enhanced_sections)

    # append architecture mermaid diagram explicitly near top
    final_md = final_md + "\n\n\n## Architecture Diagram\n\n" + ARCHITECTURE_MERMAID + "\n\n"

    # Add risk registry summary
    final_md = final_md + "## Risk Register (Top-level)\n\n"
    final_md = final_md + "| Risk | Probability | Impact | Mitigation |\n|---|---|---|---|\n"
    final_md = final_md + "| Sensitive PII leakage | Medium | High | Mask on-device; encrypt in transit & at rest |\n"
    final_md = final_md + "| RL unsafe action (destructive) | Low-Medium | High | HITL for destructive flows; sandboxed training |\n"
    final_md = final_md + "| LLM hallucination | Medium | Medium | Validation step + deterministic checks for critical actions |\n"

    return final_md, metadata


# -------------------------
# CLI
# -------------------------
def main():
    p = argparse.ArgumentParser()
    p.add_argument("--input", "-i", required=True, help="Input spec file (.md or .docx)")
    p.add_argument("--output_dir", "-o", default="./out", help="Output dir")
    p.add_argument("--project_name", default="AutoRL", help="Project name for outputs")
    p.add_argument("--author", default="AutoRL Team", help="Author for slides")
    args = p.parse_args()

    inp = Path(args.input)
    out_dir = Path(args.output_dir)
    out_dir.mkdir(parents=True, exist_ok=True)

    raw = read_input(inp)
    enhanced_md, metadata = build_enhanced_spec(raw)
    metadata["project_name"] = args.project_name

    # write enhanced md
    enhanced_path = out_dir / "spec_enhanced.md"
    enhanced_path.write_text(enhanced_md, encoding="utf-8")

    # write checklist
    checklist_path = out_dir / "spec_checklist.md"
    checklist_path.write_text(CHECKLIST_TEMPLATE, encoding="utf-8")

    # write metadata json
    metadata_path = out_dir / "spec_metadata.json"
    metadata_path.write_text(json.dumps(metadata, indent=2), encoding="utf-8")

    # create slides
    slides_path = out_dir / "slides_skeleton.pptx"
    create_slides(slides_path, title=f"{args.project_name} — Technical Spec", author=args.author)

    print("Generated:")
    print(" -", enhanced_path)
    print(" -", checklist_path)
    print(" -", metadata_path)
    print(" -", slides_path)


if __name__ == "__main__":
    main()



