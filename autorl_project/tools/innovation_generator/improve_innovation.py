
from pathlib import Path
import argparse
import json
import re
from datetime import datetime
from typing import List, Tuple, Dict
try:
    from docx import Document
    from jinja2 import Template
    from pptx import Presentation
    from pptx.util import Inches, Pt
except Exception as e:
    raise RuntimeError("Missing dependency. pip install python-docx jinja2 python-pptx") from e


# -------------------------
# Helpers
# -------------------------
def read_input(path: Path) -> str:
    if path.suffix.lower() in (".md", ".markdown", ".txt"):
        return path.read_text(encoding="utf-8")
    if path.suffix.lower() == ".docx":
        doc = Document(str(path))
        paras = [p.text for p in doc.paragraphs if p.text.strip()]
        return "\n\n".join(paras)
    raise ValueError("Unsupported input format. Use .md or .docx")


def split_sections(text: str) -> List[Tuple[str, str]]:
    """
    Split by headings (lines starting with \'#\', \'##\', or \'###\'), returning (heading, body).
    If no headings found, returns a single (\'Preamble\', full_text).
    """
    lines = text.splitlines()
    heading_re = re.compile(r\'^\\s{0,3}#{1,3}\\s+(.*)\'
)
    sections = []
    cur_h = None
    cur_lines = []
    for line in lines:
        m = heading_re.match(line)
        if m:
            # flush previous
            if cur_h is not None or cur_lines:
                sections.append((cur_h or "Preamble", "\n".join(cur_lines).strip()))
            cur_h = m.group(1).strip()
            cur_lines = []
        else:
            cur_lines.append(line)
    if cur_h is not None or cur_lines:
        sections.append((cur_h or "Preamble", "\n".join(cur_lines).strip()))
    if not sections:
        sections = [("Preamble", "")]
    return sections


# -------------------------
# Templates
# -------------------------
SECTION_TEMPLATE = Template("""### {{ title }}

{{ original | default(\'\') }}

#### Improved Executive Summary
{{ summary }}

#### Concrete Example (step-by-step)
{% for step in example %}
- {{ step }}
{% endfor %}

#### Metrics & Success Criteria
{% for metric in metrics %}
- {{ metric }}
{% endfor %}

#### Competitive Differentiation
{{ differentiation }}

#### Risks & Mitigations
| Risk | Mitigation |
|---|---|
{% for r,m in risks %}
| {{ r }} | {{ m }} |
{% endfor %}

#### Implementation Notes & Next Steps
{% for note in impl %}
- {{ note }}
{% endfor %}
""")

SLIDE_BULLETS_TEMPLATE = Template("""# Slide skeleton for \"{{ title }}\"

Bullets:
{% for b in bullets %}
- {{ b }}
{% endfor %}

Notes:
{{ notes }}
""")


# -------------------------
# Enrichment rules
# -------------------------
def enrich_section(title: str, body: str, project_name: str) -> Dict:
    t = title.lower() if title else ""
    # defaults
    summary = ""
    example = []
    metrics = []
    differentiation = ""
    risks = []
    impl = []
    bullets = []
    notes = ""

    if "executive" in t or "summary" in t or "preamble" in t:
        summary = f"{project_name} is an on-device+cloud mobile agent combining vision, LLM planning, and RL to generalize multi-step workflows across apps. It aims to reduce repetitive mobile tasks and increase task completion rate and user productivity."
        example = [
            "User: \'Pay my phone bill in the BankApp\'.",
            "Agent captures current screen, detects \'Pay\' button + amount fields via OCR/detectors.",
            "LLM planner generates sequence: TAP(\'Payments\'), TYPE(amount), TAP(\'Confirm\').",
            "HITL step for verification (if payment) — human approves.",
            "Execution agent performs actions via Appium/Accessibility; RL updates policy on success."
        ]
        metrics = [
            "Task success rate (target ≥ 95%)",
            "Average steps to completion",
            "Median end-to-end latency (target < 5s)",
            "Human interventions per 100 tasks (target < 2)"
        ]
        differentiation = "Unlike standard LLMs or scripted RPA, AutoRL learns from visual context and user feedback, enabling cross-app transfer without per-app scripting."
        risks = [("PII leakage", "Mask screenshots on-device and encrypt in transit"),
                 ("LLM hallucination onscreen", "Validate plans with deterministic checks; require HITL for destructive actions")]
        impl = ["Add sample plan JSON schema", "Include demo video (30s) showing cross-app flow", "Prepare training snapshots and evaluation scripts"]
        bullets = ["One-line elevator pitch", "Demo video + 2 short clips", "Key metrics and ask"]
        notes = "Use the above one-liner on the title slide and the demo slide."
    elif "architecture" in t or "how it works" in t:
        summary = "Show a simple dataflow: Mobile capture → Perception → Planner (LLM) → Execution → Feedback (RL + Memory). Emphasize modularity and fault-tolerance."
        example = [
            "API: POST /instruction { user_instruction, context_snapshot }",
            "Planner returns plan_json: [{action: \'tap\', target: {...}}, ...]",
            "Executor validates actions against element hit-tests before performing writes."
        ]
        metrics = ["API success rate", "Perception element detection F1 score", "Plan reuse rate"]
        differentiation = "Design separates high-level planning (LLM) from low-level execution (deterministic module), preventing LLM drift from causing destructive actions."
        risks = [("Action safety", "Executor simulates/dry-runs non-critical actions and requires HITL for destructive ones"),
                 ("Perception errors", "Combine OCR + element heuristics and fallback to semantic matching")]
        impl = ["Include sequence diagram & mermaid in spec", "Add sample API schemas and sample JSON plan"]
        bullets = ["Architecture diagram", "Sequence example", "Safety guardrails (HITL, dry-run)"]
        notes = "Add a mermaid diagram and attach sample JSON payloads in appendix."
    elif "innovation" in t or "novel" in t or "unique" in t:
        summary = f"The core innovation of {project_name} is a hybrid LLM+RL loop grounded in visual UI understanding plus a plan memory that reuses successful action chains across apps."
        example = [
            "Agent learns \'search_flow\' skill and applies it to Amazon, Maps, and Food apps with minimal fine-tuning.",
            "Plan memory retrieves prior successful sequences when UI embeddings are similar.",
            "RL policy fine-tunes at runtime using user corrections as reward signals."
        ]
        metrics = ["Transfer success: % tasks where plan from memory succeeds without changes",
                   "Learning speed: episodes-to-convergence on new app (target < 50 episodes)"]
        differentiation = "The Plan Memory + Visual Embeddings create a data moat that static script-based RPA cannot replicate."
        risks = [("Over-generalization", "Use similarity thresholds and conservative plan adaptation"),
                 ("Data privacy", "Federated updates and on-device embeddings")]
        impl = ["Implement plan memory as vector DB (Qdrant) with L2 similarity and time-decay weighting",
                "Add curriculum learning schedule for RL bootstrap"]
        bullets = ["Core innovation statement", "How plan memory works", "Transfer examples"]
        notes = "Quantify transfer experiments in appendix (small table comparing baseline scripted bot vs AutoRL)."
    elif "business" in t or "market" in t or "model" in t:
        summary = "Propose Agent-as-a-Service (AaaS), developer SDKs, and a marketplace for specialized agent skills."
        example = [
            "Freemium consumer tier: X free tasks per month.",
            "Enterprise tier: API + device fleet + data retention SLA.",
            "Marketplace: curated agent skills for travel, finance, etc."
        ]
        metrics = ["ARPU, conversion rate from free->paid", "Number of marketplace skills published"]
        differentiation = "Monetize via both subscriptions and marketplace fees; enterprise SLA + compliance is a moat."
        risks = [("Regulation/certification", "Plan for SOC2 / GDPR policies early"),
                 ("Monetization friction", "Offer trial credits and easy onboarding templates")]
        impl = ["Define pricing tiers and demo marketplace UX", "Prepare minimal SDK and sample agent"]
        bullets = ["Business model summary", "Go-to-market plan", "Pricing tiers"]
        notes = "Attach sample pricing table and MVP marketplace UX screenshots."
    elif "demo" in t or "proof" in t or "showcase" in t:
        summary = "Prepare a 90–120s demo video with 3 scenes: simple task, cross-app task, HITL safety scenario."
        example = [
            "Scene 1 (20s): \'Add event to calendar\' — agent completes without HITL.",
            "Scene 2 (40s): \'Book flight\' — agent interacts across email → airline app; shows plan and executes.",
            "Scene 3 (30s): \'Sensitive payment\' — HITL popup; human approves; agent proceeds."
        ]
        metrics = ["Demo success rate (local runs)", "Reproducibility on emulator snapshot"]
        differentiation = "Emphasize real-device visual interactions rather than API-only demos."
        risks = [("Flaky demo", "Record video of deterministic run using emulators and fixed accounts"),
                 ("Privacy in demo", "Use demo accounts and masked data")]
        impl = ["Prepare demo script, pre-record, and include timestamps", "Provide narrated subtitles for judges"]
        bullets = ["Demo flow timeline", "Key takeaways", "What to show on each slide"]
        notes = "Record both live demo and a fallback recorded video in case of live demo issues."
    else:
        # generic improvement
        summary = "Clarify the key point, add 2–3 short examples, add measurable metrics and a short demo idea."
        example = ["Add one short user scenario", "Add a short step-by-step example", "Add one metric"]
        metrics = ["At least one measurable success metric (e.g., task success rate)"]
        differentiation = "State why this is better than scripted approaches or naive LLM wrappers."
        risks = [("Underspecified", "Add measurable criteria and an explicit demo")]
        impl = ["Add short example, one demo script, and at least 2 metrics"]
        bullets = ["Clarify one-liner", "Add 1-2 examples", "Add metrics"]
        notes = "Fill gaps quickly to strengthen the submission."

    return {
        "title": title or "Preamble",
        "original": body,
        "summary": summary,
        "example": example,
        "metrics": metrics,
        "differentiation": differentiation,
        "risks": risks,
        "impl": impl,
        "bullets": bullets,
        "notes": notes
    }


# -------------------------
# PPTX helper
# -------------------------
def make_pptx(out_path: Path, project_name: str, enriched_sections: List[Dict]):
    prs = Presentation()
    # title slide
    slide_layout = prs.slide_layouts[0]
    s = prs.slides.add_slide(slide_layout)
    s.shapes.title.text = f"{project_name} — Innovation Summary"
    s.placeholders[1].text = f"Generated {datetime.utcnow().isoformat()}"

    # one slide per major enriched section (limit to 8 slides)
    for sec in enriched_sections[:8]:
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = sec["title"][:60]
        tf = slide.shapes.placeholders[1].text_frame
        for b in (sec.get("bullets") or [])[:6]:
            p = tf.add_paragraph()
            p.text = b
            p.level = 1

        # add notes for presenters
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = sec.get("notes", "")

    prs.save(str(out_path))


# -------------------------
# Main create/enhance pipeline
# -------------------------
def build_enhanced_innovation(raw_text: str, project_name: str) -> Tuple[str, Dict]:
    sections = split_sections(raw_text)
    # if no explicit headings or missing typical sections, we will add defaults
    section_titles = [s[0].lower() for s in sections]
    defaults = ["Executive Summary", "Technical Architecture", "Innovation / Novelty", "Implementation Strategy",
                "Responsible AI & Safety", "Demo Plan", "Business Model", "Roadmap & Next Steps"]
    # Promote existing: if a default missing, insert with empty body
    existing_titles = {s[0].lower(): s[1] for s in sections}
    ordered_sections = []
    for d in defaults:
        key = d.lower()
        ordered_sections.append((d, existing_titles.get(key, "")))

    enriched = []
    metadata_sections = []
    for title, body in ordered_sections:
        info = enrich_section(title, body, project_name)
        enriched.append(info)
        metadata_sections.append({"title": title, "has_original": bool(body)})

    # assemble markdown output
    header = f"# {project_name} — Innovation (Enhanced)\n\n_Generated: {datetime.utcnow().isoformat()}_\n\n"
    md_parts = [header]
    for sec in enriched:
        md = SECTION_TEMPLATE.render(
            title=sec["title"],
            original=sec["original"],
            summary=sec["summary"],
            example=sec["example"],
            metrics=sec["metrics"],
            differentiation=sec["differentiation"],
            risks=sec["risks"],
            impl=sec["impl"]
        )
        md_parts.append(md)
    full_md = "\n\n".join(md_parts)

    # short checklist
    checklist = [
        "- [ ] One-sentence innovation statement",
        "- [ ] 2 concrete user scenarios with step-by-step flows",
        "- [ ] Measurable metrics (task success rate, latency, human interventions)",
        "- [ ] Demo script + recorded fallback",
        "- [ ] Risk register + mitigations (PII, hallucination, destructive actions)",
        "- [ ] Slides (title, demo, architecture, metrics, ask)"
    ]
    metadata = {"generated_at": datetime.utcnow().isoformat(), "project": project_name, "sections": metadata_sections, "checklist": checklist}
    return full_md, metadata, enriched


# -------------------------
# CLI
# -------------------------
def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--input", "-i", required=True, help="Input innovation file (.md or .docx)")
    parser.add_argument("--output_dir", "-o", default="./out", help="Output directory")
    parser.add_argument("--project", "-p", default="AutoRL", help="Project name")
    args = parser.parse_args()

    inp = Path(args.input)
    out_dir = Path(args.output_dir)
    out_dir.mkdir(parents=True, exist_ok=True)

    raw = read_input(inp)
    enhanced_md, metadata, enriched = build_enhanced_innovation(raw, args.project)

    md_path = out_dir / "innovation_enhanced.md"
    md_path.write_text(enhanced_md, encoding="utf-8")

    meta_path = out_dir / "innovation_metadata.json"
    meta_path.write_text(json.dumps(metadata, indent=2), encoding="utf-8")

    checklist_path = out_dir / "innovation_checklist.md"
    checklist_path.write_text("# Innovation Checklist\n\n" + "\n".join(metadata["checklist"]), encoding="utf-8")

    pptx_path = out_dir / "innovation_slides.pptx"
    make_pptx(pptx_path, args.project, enriched)

    print("Generated:")
    print(" -", md_path)
    print(" -", meta_path)
    print(" -", checklist_path)
    print(" -", pptx_path)


if __name__ == "__main__":
    main()

